"""Extracción de texto de documentos en varios formatos.

Punto de entrada: ``extraer_texto(nombre_archivo, datos)`` -> (texto, tipo, truncado).
Cada formato tiene su propio extractor; el dispatcher elige por la extensión.
"""
from __future__ import annotations

import csv
import io

# Tope de caracteres que se conservan; el resto se descarta antes de clasificar
# (el modelo solo necesita una muestra representativa del documento).
MAX_CARACTERES = 100_000


class FormatoNoSoportado(Exception):
    """El formato del archivo no tiene un extractor disponible."""


def _extraer_pdf(datos: bytes) -> str:
    from pypdf import PdfReader

    lector = PdfReader(io.BytesIO(datos))
    return "\n".join((pagina.extract_text() or "") for pagina in lector.pages)


def _extraer_docx(datos: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(datos))
    partes = [parrafo.text for parrafo in doc.paragraphs]
    # También el texto contenido en tablas.
    for tabla in doc.tables:
        for fila in tabla.rows:
            for celda in fila.cells:
                partes.append(celda.text)
    return "\n".join(partes)


def _decodificar(datos: bytes) -> str:
    """Decodifica bytes a texto detectando automáticamente la codificación."""
    from charset_normalizer import from_bytes

    resultado = from_bytes(datos).best()
    if resultado is not None:
        return str(resultado)
    return datos.decode("utf-8", errors="replace")


def _extraer_txt(datos: bytes) -> str:
    return _decodificar(datos)


def _extraer_csv(datos: bytes) -> str:
    texto = _decodificar(datos)
    celdas: list[str] = []
    for fila in csv.reader(io.StringIO(texto)):
        celdas.extend(celda for celda in fila if celda)
    return " ".join(celdas)


def _extraer_xlsx(datos: bytes) -> str:
    from openpyxl import load_workbook

    libro = load_workbook(io.BytesIO(datos), read_only=True, data_only=True)
    celdas: list[str] = []
    for hoja in libro.worksheets:
        for fila in hoja.iter_rows(values_only=True):
            for celda in fila:
                if celda is not None:
                    celdas.append(str(celda))
    libro.close()
    return " ".join(celdas)


def _extraer_pptx(datos: bytes) -> str:
    from pptx import Presentation

    presentacion = Presentation(io.BytesIO(datos))
    partes: list[str] = []
    for diapositiva in presentacion.slides:
        for forma in diapositiva.shapes:
            if forma.has_text_frame:
                partes.append(forma.text_frame.text)
    return "\n".join(partes)


# Dispatcher: extensión -> función extractora.
_EXTRACTORES = {
    "pdf": _extraer_pdf,
    "docx": _extraer_docx,
    "txt": _extraer_txt,
    "csv": _extraer_csv,
    "xlsx": _extraer_xlsx,
    "pptx": _extraer_pptx,
}

FORMATOS_SOPORTADOS: list[str] = sorted(_EXTRACTORES)


def extraer_texto(nombre_archivo: str, datos: bytes) -> tuple[str, str, bool]:
    """Extrae el texto de un documento.

    Devuelve ``(texto, tipo_archivo, truncado)``. Lanza ``FormatoNoSoportado``
    si la extensión del archivo no tiene un extractor disponible.
    """
    extension = ""
    if "." in nombre_archivo:
        extension = nombre_archivo.rsplit(".", 1)[-1].lower()

    extractor = _EXTRACTORES.get(extension)
    if extractor is None:
        raise FormatoNoSoportado(
            f"Formato '.{extension}' no soportado. "
            f"Soportados: {', '.join(FORMATOS_SOPORTADOS)}."
        )

    texto = (extractor(datos) or "").strip()
    truncado = len(texto) > MAX_CARACTERES
    if truncado:
        texto = texto[:MAX_CARACTERES]
    return texto, extension, truncado
