"""Pruebas de extracción de texto. Genera un fichero de cada formato al vuelo."""
import io

import pytest

import extraccion
from extraccion import FormatoNoSoportado, extraer_texto

TEXTO = "Informacion confidencial: el DNI es 12345678Z."


def test_txt():
    texto, tipo, truncado, advertencias = extraer_texto("doc.txt", TEXTO.encode("utf-8"))
    assert "confidencial" in texto
    assert tipo == "txt"
    assert truncado is False
    assert advertencias == []


def test_csv():
    datos = "nombre,dni\nAna,12345678Z\n".encode("utf-8")
    texto, tipo, _, _ = extraer_texto("doc.csv", datos)
    assert "12345678Z" in texto
    assert tipo == "csv"


def test_docx():
    from docx import Document

    doc = Document()
    doc.add_paragraph(TEXTO)
    buffer = io.BytesIO()
    doc.save(buffer)
    texto, tipo, _, _ = extraer_texto("doc.docx", buffer.getvalue())
    assert "confidencial" in texto
    assert tipo == "docx"


def test_xlsx():
    from openpyxl import Workbook

    libro = Workbook()
    libro.active["A1"] = "DNI"
    libro.active["B1"] = "12345678Z"
    buffer = io.BytesIO()
    libro.save(buffer)
    texto, tipo, _, _ = extraer_texto("doc.xlsx", buffer.getvalue())
    assert "12345678Z" in texto
    assert tipo == "xlsx"


def test_pptx():
    from pptx import Presentation
    from pptx.util import Inches

    pres = Presentation()
    diapositiva = pres.slides.add_slide(pres.slide_layouts[6])  # diapositiva en blanco
    cuadro = diapositiva.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    cuadro.text_frame.text = TEXTO
    buffer = io.BytesIO()
    pres.save(buffer)
    texto, tipo, _, _ = extraer_texto("doc.pptx", buffer.getvalue())
    assert "confidencial" in texto
    assert tipo == "pptx"


def test_pdf():
    from reportlab.pdfgen import canvas

    buffer = io.BytesIO()
    lienzo = canvas.Canvas(buffer)
    lienzo.drawString(100, 700, TEXTO)
    lienzo.save()
    texto, tipo, _, _ = extraer_texto("doc.pdf", buffer.getvalue())
    assert "confidencial" in texto
    assert tipo == "pdf"


def test_audio_usa_whisper_antes_de_markitdown(monkeypatch):
    monkeypatch.setattr(extraccion, "_transcribir_audio_whisper", lambda _ruta: "texto de audio")
    texto, tipo, truncado, advertencias = extraer_texto("audio.mpeg", b"datos")
    assert texto == "texto de audio"
    assert tipo == "mpeg"
    assert truncado is False
    assert advertencias == []


def test_formato_no_soportado():
    with pytest.raises(FormatoNoSoportado):
        extraer_texto("imagen.bmp", b"datos")
